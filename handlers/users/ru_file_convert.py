
from loader import dp, bot
import os
import pandas as pd
from aiogram import Bot, Dispatcher, types
from PyPDF2 import PdfReader
from pptx import Presentation
from PIL import Image
from gtts import gTTS
import zipfile  # ZIP fayllarni boshqarish uchun kutubxona


# Callback handler for PDF to PPTX conversion
@dp.callback_query_handler(lambda c: c.data == "ru_pdf_to_pptx")
async def convert_pdf_to_pptx_handler(callback_query: types.CallbackQuery):
    await bot.send_message(callback_query.from_user.id, "–û—Ç–ø—Ä–∞–≤—å—Ç–µ PDF-—Ñ–∞–π–ª! üìÑ")
    await bot.answer_callback_query(callback_query.id)

# Faylni PDF dan PPTX ga o'zgartirish
@dp.message_handler(content_types=types.ContentType.DOCUMENT)
async def handle_pdf_to_pptx(message: types.Message):
    if message.document.mime_type == "application/pdf":
        file_path = f"/tmp/{message.document.file_name}"
        await message.document.download(file_path)

        pptx_path = convert_pdf_to_pptx(file_path)  # Konversiya funksiyasini chaqirish
        await message.reply_document(types.InputFile(pptx_path), caption="PDF-—Ñ–∞–π–ª –∫–æ–Ω–≤–µ—Ä—Ç–∏—Ä–æ–≤–∞–Ω –≤ PPTX.! üìä")
        os.remove(file_path)
        os.remove(pptx_path)


# Faylni PDF dan PPTX ga o'zgartirish
def convert_pdf_to_pptx(pdf_path):
    reader = PdfReader(pdf_path)
    presentation = Presentation()

    for page in reader.pages:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        text = page.extract_text()
        textbox = slide.shapes.add_textbox(
            left=0, top=0, width=presentation.slide_width, height=presentation.slide_height
        )
        textbox.text = text

    pptx_path = pdf_path.replace('.pdf', '.pptx')
    presentation.save(pptx_path)
    return pptx_path




# JPG ni PNG ga ozgartirish

# Callback handler for JPG to PNG conversion
@dp.callback_query_handler(lambda c: c.data == "ru_jpg_to_png")
async def convert_jpg_to_png_handler(callback_query: types.CallbackQuery):
    await bot.send_message(callback_query.from_user.id, "–û—Ç–ø—Ä–∞–≤—å—Ç–µ —Å–≤–æ–π JPG-—Ñ–∞–π–ª! üñºÔ∏è")
    await bot.answer_callback_query(callback_query.id)


# Faylni JPG dan PNG ga o'zgartirish
@dp.message_handler(content_types=types.ContentType.PHOTO)
async def handle_jpg_to_png(message: types.Message):
    photo = message.photo[-1]  # Yuborilgan fotosuratning eng katta versiyasini olish
    file = await photo.get_file()
    file_path = f"/tmp/{file.file_id}.jpg"
    await file.download(file_path)

    png_path = convert_jpg_to_png(file_path)  # Konversiya funksiyasini chaqirish
    await message.reply_document(types.InputFile(png_path), caption="JPG-—Ñ–∞–π–ª –∫–æ–Ω–≤–µ—Ä—Ç–∏—Ä–æ–≤–∞–Ω –≤ PNG! üì∑")
    os.remove(file_path)
    os.remove(png_path)


# JPG ni PNG ga aylantirish
def convert_jpg_to_png(jpg_path):
    img = Image.open(jpg_path)
    png_path = jpg_path.replace('.jpg', '.png')
    img.save(png_path, 'PNG')
    return png_path




# JSON ni CSV ga o'zgartirish

# Callback handler for JSON to CSV conversion
@dp.callback_query_handler(lambda c: c.data == "ru_json_csv")
async def convert_json_to_csv_handler(callback_query: types.CallbackQuery):
    await bot.send_message(callback_query.from_user.id, "–û—Ç–ø—Ä–∞–≤—å—Ç–µ —Å–≤–æ–π JSON-—Ñ–∞–π–ª! üìÑ")
    await bot.answer_callback_query(callback_query.id)

# Faylni JSON dan CSV ga o'zgartirish
@dp.message_handler(content_types=types.ContentType.DOCUMENT)
async def handle_json_to_csv(message: types.Message):
    if message.document.mime_type == "application/json":
        file_path = f"/tmp/{message.document.file_name}"
        await message.document.download(file_path)

        csv_path = convert_json_to_csv(file_path)  # Konversiya funksiyasini chaqirish
        await message.reply_document(types.InputFile(csv_path), caption="–§–∞–π–ª JSON –ø—Ä–µ–æ–±—Ä–∞–∑–æ–≤–∞–Ω –≤ CSV! üìà")
        os.remove(file_path)
        os.remove(csv_path)

# JSON ni CSV ga aylantirish
def convert_json_to_csv(json_path):
    data = pd.read_json(json_path)
    csv_path = json_path.replace('.json', '.csv')
    data.to_csv(csv_path, index=False)
    return csv_path



# Fayllarni ZIPga o'zgartirish

# Callback handler for ZIP creation
@dp.callback_query_handler(lambda c: c.data == "ru_zip")
async def zip_prompt(callback_query: types.CallbackQuery):
    await bot.send_message(callback_query.from_user.id, "–û—Ç–ø—Ä–∞–≤—å—Ç–µ —Å–≤–æ–∏ —Ñ–∞–π–ª—ã –¥–ª—è —Å–æ–∑–¥–∞–Ω–∏—è ZIP-–∞—Ä—Ö–∏–≤–∞! üì¶")
    await bot.answer_callback_query(callback_query.id)

# Faylni ZIPga qo'shish
@dp.message_handler(content_types=types.ContentType.DOCUMENT)
async def handle_document_zip(message: types.Message):
    document = message.document
    file_path = f"/tmp/{document.file_name}"
    await document.download(file_path)

    zip_path = "/tmp/archive.zip"
    try:
        with zipfile.ZipFile(zip_path, 'a') as archive:
            archive.write(file_path, arcname=document.file_name)  # Faylni ZIP ichiga joylash
        await message.reply("–§–∞–π–ª –¥–æ–±–∞–≤–ª–µ–Ω –≤ ZIP-–∞—Ä—Ö–∏–≤. –û—Ç–ø—Ä–∞–≤—å—Ç–µ /finish_zip, —á—Ç–æ–±—ã —Ä–∞—Å–ø–∞–∫–æ–≤–∞—Ç—å –∞—Ä—Ö–∏–≤.. üì¶")
        os.remove(file_path)  # Ortiqcha joyni tozalash
    except Exception as e:
        await message.reply(f"–ü—Ä–æ–∏–∑–æ—à–ª–∞ –æ—à–∏–±–∫–∞ –ø—Ä–∏ —Å–æ–∑–¥–∞–Ω–∏–∏ ZIP-–∞—Ä—Ö–∏–≤–∞.: {e}")

# ZIP arxivni tugatish va foydalanuvchiga yuborish
@dp.message_handler(commands=['finish_zip'])
async def finish_zip(message: types.Message):
    zip_path = "/tmp/archive.zip"
    if os.path.exists(zip_path):
        await message.reply_document(types.InputFile(zip_path), caption="ZIP-–∞—Ä—Ö–∏–≤ –≥–æ—Ç–æ–≤.! üì¶")
        os.remove(zip_path)  # ZIPni foydalanuvchiga yuborganingizdan keyin o'chirish
    else:
        await message.reply("ZIP-–∞—Ä—Ö–∏–≤ –µ—â–µ –Ω–µ —Å–æ–∑–¥–∞–Ω. –û—Ç–ø—Ä–∞–≤–∏—Ç—å —Ñ–∞–π–ª –∏ –¥–æ–±–∞–≤–∏—Ç—å –≤ ZIP! üì¶")
